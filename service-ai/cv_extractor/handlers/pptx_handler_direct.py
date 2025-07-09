import os
from pptx import Presentation
from cv_extractor.utils.file_utils import save_text_as_txt, log_info, log_warn, log_error
from cv_extractor.utils.parser import clean_text
from cv_extractor.utils.structuring_api_gemini_direct import structure_with_gemini_direct
from cv_extractor.utils.structuring_api_direct import structure_with_openai_direct
from cv_extractor.utils.structuring_common import save_structured_result

def handle_pptx_direct(file_path):
    filename = os.path.basename(file_path)
    log_info(f"🔐 Handling PPTX with direct structuring: {filename}")

    try:
        # === Extract text from slides ===
        prs = Presentation(file_path)
        text_chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_chunks.append(shape.text.strip())
        text = "\n".join(text_chunks)

        if not text.strip():
            log_warn(f"⚠️ No text extracted from {filename}")
            return None

        # === Clean & Save cleaned text ===
        folder = "output_pptx"
        os.makedirs(folder, exist_ok=True)
        cleaned_text = clean_text(text)
        out_path = os.path.join(folder, os.path.splitext(filename)[0] + "_text.txt")
        save_text_as_txt(cleaned_text, out_path)
        log_info(f"✅ Cleaned text saved to: {out_path}")

        # === Try Gemini Direct ===
        try:
            structured_raw = structure_with_gemini_direct(cleaned_text)
            log_info("✅ Structured using Gemini Direct")
        except Exception as gemini_err:
            log_error(f"⚠️ Gemini Direct failed: {gemini_err}")
            try:
                structured_raw = structure_with_openai_direct(cleaned_text)
                log_info("✅ Structured using OpenAI Direct fallback")
            except Exception as openai_err:
                log_error(f"❌ OpenAI Direct also failed: {openai_err}")
                return None

        # === Save structured result and return JSON ===
        structured_json = save_structured_result(file_path, structured_raw)
        return structured_json

    except Exception as e:
        log_error(f"❌ Fatal error while handling PPTX file {filename}: {e}")
        return None
